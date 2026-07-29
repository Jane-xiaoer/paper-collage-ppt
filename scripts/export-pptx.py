#!/usr/bin/env python3
"""
剪纸拼贴 HTML deck → PPTX 导出。

原理:无头 Chrome 逐页截图(1920×1080),再拼成 16:9 的 PowerPoint——
每页是一张整页图片,视觉 100% 还原(含当前选定配色);代价是文字在 PPTX
里不可编辑、入场动画丢失。要改内容回 HTML 改完重新导。

用法:
    python3 export-pptx.py deck.html                    # 输出 deck.pptx
    python3 export-pptx.py deck.html -o 分享.pptx
    python3 export-pptx.py deck.html --theme t13        # 指定配色导出(不带则用文件里的初始主题)

依赖:Google Chrome + python-pptx(缺了会提示安装命令)。
"""
import argparse
import re
import subprocess
import sys
import tempfile
from pathlib import Path

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("-o", "--out")
    ap.add_argument("--theme", help="tN 配色编号,如 t13")
    ap.add_argument("--size", default="1920x1080")
    args = ap.parse_args()

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        sys.exit("ERROR: 缺 python-pptx。安装:python3 -m pip install --user python-pptx")

    deck = Path(args.deck).expanduser().resolve()
    if not deck.exists():
        sys.exit(f"ERROR: 找不到 {deck}")
    if not Path(CHROME).exists():
        sys.exit("ERROR: 找不到 Google Chrome")
    w, h = (int(x) for x in args.size.split("x"))

    html = deck.read_text()
    n = html.count('<section class="slide')
    if n == 0:
        sys.exit("ERROR: 文件里没有 slide")

    src = deck
    if args.theme:
        if not re.search(rf'body\.{args.theme}\b', html):
            sys.exit(f"ERROR: 模板里没有主题 {args.theme}")
        tmp = Path(tempfile.gettempdir()) / f"pcx-{deck.stem}.html"
        # 覆写初始主题,并清掉 localStorage 干扰(改 KEY 让它读不到旧选择)
        patched = re.sub(r'<body class="t\d+"', f'<body class="{args.theme}"', html, count=1)
        patched = patched.replace("localStorage.getItem(KEY)", "null")
        tmp.write_text(patched)
        src = tmp

    print(f"  {n} 页,逐页截图({w}x{h})...")
    shots = []
    with tempfile.TemporaryDirectory() as td:
        for i in range(1, n + 1):
            png = Path(td) / f"s{i:02d}.png"
            subprocess.run(
                [CHROME, "--headless", "--disable-gpu", f"--window-size={w},{h}",
                 "--virtual-time-budget=6000", "--hide-scrollbars",
                 f"--screenshot={png}", f"file://{src}#{i}"],
                check=True, capture_output=True,
            )
            if not png.exists():
                sys.exit(f"ERROR: 第 {i} 页截图失败")
            shots.append(png)
            print(f"  ✓ {i}/{n}")

        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 16:9
        prs.slide_height = Emu(6858000)
        blank = prs.slide_layouts[6]
        for png in shots:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(png), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        out = Path(args.out).expanduser() if args.out else deck.with_suffix(".pptx")
        prs.save(out)

    print(f"  ✓ 导出完成: {out}  ({out.stat().st_size // 1024}KB)")


if __name__ == "__main__":
    main()
