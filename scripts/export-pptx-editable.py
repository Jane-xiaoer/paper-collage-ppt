#!/usr/bin/env python3
"""
剪纸拼贴 HTML deck → 可编辑 PPTX(三层分解导出)。

分解方式:
  层1 背景  = 底纹/山景/纸屑等装饰(逐页截图,不含纸片和文字)
  层2 纸片  = 每张纸片/撕纸条/标签/剪贴字底片 = 独立图片形状(canvas 实时
             套色渲染,位置/大小/旋转对位)——可选中、可拖动、可拉伸
  层3 文字  = 真文本框(位置/字号/颜色/旋转/字距对位)——直接改字

文字变长 → 点住纸片拉大即可,纸片纹理跟着拉伸(卡纸纹理拉伸不易穿帮)。

用法:
    python3 export-pptx-editable.py deck.html                    # deck-editable.pptx
    python3 export-pptx-editable.py deck.html -o 分享.pptx --theme t15

边界(诚实声明):
- 改字/挪纸片/拉纸片 ✅;换配色/改布局结构 ❌(回 HTML 重导)
- 只有「图片材质」的纸片会被拆出来(卡纸版全部是);纯 CSS 色块纸片
  (轻量版部分元素)保持烤在背景里
- 马克笔垫色(uline)烤在背景里,不跟随文字变长

依赖:Google Chrome + python-pptx(+ PIL 可选,用于压缩)。
"""
import argparse
import base64
import json
import re
import subprocess
import sys
import tempfile
from pathlib import Path

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
VW, VH = 1600, 900
EMU_PER_PX = 7620          # 12192000 EMU / 1600 px
PT_PER_PX = 72 / 120

FORCE_FINAL = """
<style id="__pcx">
#skinbar,#dots,#pageno{display:none!important}
[data-drop]{opacity:1!important;transform:translateY(0) rotate(var(--rot,0deg)) scale(1)!important;transition:none!important}
*{transition:none!important;animation:none!important}
</style>
"""

# 共用采集逻辑:mode='extract' 输出 JSON;mode='mask' 隐藏已导出元素
COLLECT_JS = r"""
async function __pcxRun(mode){
  const PIECE_SEL = '.piece,.torn-strip,.torn-b,.label,.chips > span,.tape';
  function angleOf(el, stopAt){
    let a = 0, n = el;
    while(n && n !== stopAt && n.nodeType === 1){
      const tr = getComputedStyle(n).transform;
      if(tr && tr !== 'none'){
        const m = tr.match(/matrix\(([-\d.e]+),\s*([-\d.e]+)/);
        if(m) a += Math.atan2(parseFloat(m[2]), parseFloat(m[1])) * 180 / Math.PI;
      }
      n = n.parentElement;
    }
    return a;
  }
  function hasDirectText(el){
    for(const c of el.childNodes)
      if(c.nodeType === 3 && c.textContent.trim()) return true;
    return false;
  }
  function bgUrl(cs){
    const m = cs.backgroundImage.match(/url\("?(data:[^")]+)"?\)/);
    return m ? m[1] : null;
  }
  const imgCache = {};
  function loadImg(src){
    if(imgCache[src]) return imgCache[src];
    imgCache[src] = new Promise(res => {
      const im = new Image();
      im.onload = () => res(im); im.onerror = () => res(null);
      im.src = src;
    });
    return imgCache[src];
  }
  async function renderPiece(el, cs){
    const src = bgUrl(cs);
    if(!src) return null;                       // 纯 CSS 色块:不拆,留背景
    const img = await loadImg(src);
    if(!img) return null;
    const w = Math.max(el.offsetWidth, 4) * 2, h = Math.max(el.offsetHeight, 4) * 2;
    const cv = document.createElement('canvas');
    cv.width = w; cv.height = h;
    const ctx = cv.getContext('2d');
    ctx.globalAlpha = parseFloat(cs.opacity) || 1;
    ctx.drawImage(img, 0, 0, w, h);
    const bc = cs.backgroundColor;
    const rgba = bc.match(/[\d.]+/g);
    const hasTint = rgba && (rgba.length < 4 || parseFloat(rgba[3]) > 0.01);
    if(hasTint && cs.backgroundBlendMode.includes('multiply')){
      ctx.globalCompositeOperation = 'multiply';
      ctx.fillStyle = bc; ctx.fillRect(0, 0, w, h);
      ctx.globalCompositeOperation = 'destination-in';
      ctx.drawImage(img, 0, 0, w, h);
    }
    return cv.toDataURL('image/png').split(',')[1];
  }
  const texts = [], pieces = [];
  const slides = [...document.querySelectorAll('#deck > .slide')];
  for(let si = 0; si < slides.length; si++){
    const slide = slides[si];
    const srect = slide.getBoundingClientRect();
    // —— 纸片层 ——
    for(const el of slide.querySelectorAll(PIECE_SEL)){
      const cs = getComputedStyle(el);
      const r = el.getBoundingClientRect();
      if(r.width < 3 || r.height < 3) continue;
      const png = (mode === 'extract') ? await renderPiece(el, cs) : (bgUrl(cs) ? 'x' : null);
      if(!png) continue;                        // 拆不动的留在背景里
      if(mode === 'mask'){
        el.style.setProperty('background', 'none', 'important');
        el.style.setProperty('-webkit-mask-image', 'none', 'important');
        el.style.setProperty('mask-image', 'none', 'important');
        el.style.setProperty('border', 'none', 'important');
        el.style.setProperty('filter', 'none', 'important');
        el.style.setProperty('box-shadow', 'none', 'important');
      } else {
        pieces.push({slide: si, png: png,
          cx: r.left + r.width/2 - srect.left, cy: r.top + r.height/2 - srect.top,
          w: el.offsetWidth, h: el.offsetHeight, rot: angleOf(el, slide.parentElement)});
      }
    }
    // —— 文字层 ——
    const collected = [];
    const isIn = el => collected.some(c => c === el || c.contains(el));
    for(const el of slide.querySelectorAll('*')){
      if(isIn(el)) continue;
      if(['SCRIPT','STYLE','SVG','IMG'].includes(el.tagName)) continue;
      let targets = null;
      if(el.classList.contains('chips')){
        targets = [...el.children].filter(c => c.tagName === 'SPAN' && c.innerText.trim());
        collected.push(el);
      } else if(hasDirectText(el)){
        targets = [el]; collected.push(el);
      }
      if(!targets) continue;
      for(const t of targets){
        const r = t.getBoundingClientRect();
        if(r.width < 2 || r.height < 2) continue;
        const cs = getComputedStyle(t);
        if(mode === 'mask'){
          t.style.setProperty('color','transparent','important');
          t.style.setProperty('-webkit-text-fill-color','transparent','important');
          t.querySelectorAll('*').forEach(d => {
            d.style.setProperty('color','transparent','important');
            d.style.setProperty('-webkit-text-fill-color','transparent','important');
          });
          continue;
        }
        const pad = ['Top','Right','Bottom','Left'].map(s => parseFloat(cs['padding'+s]) || 0);
        texts.push({
          slide: si, text: t.innerText,
          cx: r.left + r.width/2 - srect.left, cy: r.top + r.height/2 - srect.top,
          w: t.offsetWidth, h: t.offsetHeight,
          rot: angleOf(t, slide.parentElement),
          size: parseFloat(cs.fontSize),
          lh: parseFloat(cs.lineHeight) || parseFloat(cs.fontSize) * 1.2,
          weight: parseInt(cs.fontWeight) || 400,
          color: cs.color, align: cs.textAlign,
          family: cs.fontFamily.split(',')[0].replace(/["']/g,'').trim(),
          spacing: parseFloat(cs.letterSpacing) || 0, pad: pad,
        });
      }
    }
  }
  if(mode === 'extract'){
    const pre = document.createElement('pre');
    pre.id = '__pcxout';
    pre.textContent = btoa(unescape(encodeURIComponent(JSON.stringify({texts, pieces}))));
    document.body.appendChild(pre);
    document.title = 'PCXREADY';
  } else {
    document.title = 'PCXMASKED';
  }
}
"""

EXTRACT_PAGE = COLLECT_JS + "document.fonts.ready.then(()=>setTimeout(()=>__pcxRun('extract'),300));"
MASK_PAGE = COLLECT_JS + "document.fonts.ready.then(()=>setTimeout(()=>__pcxRun('mask'),300));"

FONT_MAP = {
    "Noto Sans SC": "PingFang SC", "-apple-system": "PingFang SC",
    "Noto Serif SC": "Songti SC", "Ma Shan Zheng": "Kaiti SC",
    "STSongti-SC-Black": "Songti SC",
}


def run_chrome(args, timeout=240):
    return subprocess.run([CHROME, "--headless", "--disable-gpu", *args],
                          capture_output=True, text=False, timeout=timeout)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("-o", "--out")
    ap.add_argument("--theme", help="tN 配色编号")
    args = ap.parse_args()

    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
    except ImportError:
        sys.exit("ERROR: 缺 python-pptx。安装:python3 -m pip install --user python-pptx")

    deck = Path(args.deck).expanduser().resolve()
    if not deck.exists():
        sys.exit(f"ERROR: 找不到 {deck}")
    html = deck.read_text()
    n = html.count('<section class="slide')
    if args.theme:
        html = re.sub(r'<body class="t\d+"', f'<body class="{args.theme}"', html, count=1)
        html = html.replace("localStorage.getItem(KEY)", "null")

    td = Path(tempfile.mkdtemp(prefix="pcx-"))

    # Pass 1: 提取文字几何 + 纸片渲染
    (td/"extract.html").write_text(html.replace("</body>", FORCE_FINAL + "<script>" + EXTRACT_PAGE + "</script></body>"))
    r = run_chrome([f"--window-size={VW},{VH}", "--virtual-time-budget=20000",
                    "--dump-dom", f"file://{td}/extract.html"])
    m = re.search(rb'id="__pcxout">([A-Za-z0-9+/=]+)<', r.stdout)
    if not m:
        sys.exit("ERROR: 提取失败(dump-dom 无数据)")
    data = json.loads(base64.b64decode(m.group(1)).decode("utf-8"))
    texts, pieces = data["texts"], data["pieces"]
    print(f"  ✓ 提取 {len(texts)} 个文本框 + {len(pieces)} 张独立纸片 / {n} 页")

    # Pass 2: 纸片+文字隐藏后逐页截背景
    (td/"mask.html").write_text(html.replace("</body>", FORCE_FINAL + "<script>" + MASK_PAGE + "</script></body>"))
    shots = []
    for i in range(1, n+1):
        png = td/f"s{i:02d}.png"
        run_chrome([f"--window-size={VW},{VH}", "--force-device-scale-factor=2",
                    "--virtual-time-budget=12000", "--hide-scrollbars",
                    f"--screenshot={png}", f"file://{td}/mask.html#{i}"])
        if not png.exists():
            sys.exit(f"ERROR: 第 {i} 页背景截图失败")
        try:
            from PIL import Image
            jpg = td/f"s{i:02d}.jpg"
            Image.open(png).convert("RGB").save(jpg, "JPEG", quality=82)
            png.unlink(); png = jpg
        except ImportError:
            pass
        shots.append(png)
        print(f"  ✓ 背景 {i}/{n}")

    # 纸片 PNG 落盘(可选压缩)
    piece_files = []
    for idx, pc in enumerate(pieces):
        f = td/f"p{idx:03d}.png"
        f.write_bytes(base64.b64decode(pc["png"]))
        try:
            from PIL import Image
            im = Image.open(f)
            if im.width > 1200:
                im.thumbnail((1200, 1200), Image.LANCZOS)
                im.save(f, optimize=True)
        except ImportError:
            pass
        piece_files.append(f)

    # Pass 3: 组装 PPTX
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    ALIGN = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}

    for si in range(n):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(shots[si]), 0, 0, width=prs.slide_width, height=prs.slide_height)
        for idx, pc in enumerate(pieces):
            if pc["slide"] != si: continue
            w, h = pc["w"]*EMU_PER_PX, pc["h"]*EMU_PER_PX
            pic = slide.shapes.add_picture(str(piece_files[idx]),
                Emu(int(pc["cx"]*EMU_PER_PX - w/2)), Emu(int(pc["cy"]*EMU_PER_PX - h/2)),
                width=Emu(int(w)), height=Emu(int(h)))
            pic.rotation = pc["rot"]
        for it in (x for x in texts if x["slide"] == si):
            w, h = it["w"]*EMU_PER_PX*1.08 + 4*EMU_PER_PX, it["h"]*EMU_PER_PX
            left, top = it["cx"]*EMU_PER_PX - w/2, it["cy"]*EMU_PER_PX - h/2
            box = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(w)), Emu(int(h)))
            box.rotation = it["rot"]
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Emu(int(it["pad"][3]*EMU_PER_PX))
            tf.margin_right = Emu(int(it["pad"][1]*EMU_PER_PX))
            tf.margin_top = Emu(int(it["pad"][0]*EMU_PER_PX))
            tf.margin_bottom = Emu(int(it["pad"][2]*EMU_PER_PX))
            rgb = re.findall(r"\d+", it["color"])[:3]
            color = RGBColor(*(int(v) for v in rgb)) if len(rgb) == 3 else RGBColor(0,0,0)
            fam = FONT_MAP.get(it["family"], it["family"])
            for li, line in enumerate(it["text"].split("\n")):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = ALIGN.get(it["align"], PP_ALIGN.LEFT)
                p.line_spacing = max(it["lh"]/it["size"], 0.9)
                run = p.add_run(); run.text = line
                f = run.font
                f.size = Pt(round(it["size"]*PT_PER_PX, 1))
                f.bold = it["weight"] >= 600
                f.color.rgb = color
                f.name = fam
                rPr = run._r.get_or_add_rPr()
                ea = rPr.find(qn("a:ea"))
                if ea is None:
                    ea = rPr.makeelement(qn("a:ea"), {}); rPr.append(ea)
                ea.set("typeface", fam)
                if it["spacing"]:
                    rPr.set("spc", str(int(it["spacing"]*PT_PER_PX*100)))

    out = Path(args.out).expanduser() if args.out else deck.with_name(deck.stem + "-editable.pptx")
    prs.save(out)
    print(f"  ✓ 可编辑 PPTX: {out}  ({out.stat().st_size//1024}KB;{len(pieces)} 张纸片图形 + {len(texts)} 个文本框,全部可选中编辑)")


if __name__ == "__main__":
    main()
