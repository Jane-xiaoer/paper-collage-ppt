#!/usr/bin/env python3
"""Export Haoqi 3D HTML deck to a layered, editable PPTX.

Layers per slide:
  1. WebGL environment, lighting and unregistered scenery (one background image)
  2. Registered Three.js objects (independent transparent picture objects)
  3. CSS panels, guide lines and swatches (native PowerPoint shapes)
  4. Clay/sticker PNGs (independent movable picture objects)
  5. DOM copy (real editable PowerPoint text boxes)

The deck runtime reports its own page count and registered 3D objects through
window.__PPTX_EXPORT__; no slide count or object identity is hard-coded here.

Usage:
  python3 scripts/export-haoqi-pptx-editable.py styles/haoqi-3d/demo.html \
      -o haoqi-3d-editable.pptx --scheme plus
"""
import argparse, base64, contextlib, http.server, json, re, socketserver
import subprocess, tempfile, threading, urllib.parse
from pathlib import Path

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
VW, VH = 1600, 900
EMU_PER_PX = 7620
PT_PER_PX = 72 / 120

COLLECT = r"""
<script>
(async()=>{
 await document.fonts.ready; await new Promise(r=>setTimeout(r,1200));
 const root=document.documentElement, texts=[], images=[], shapes=[];
 const angle=el=>{let a=0,n=el;while(n&&n!==document.body){const t=getComputedStyle(n).transform;
  const m=t&&t.match(/matrix\(([-\d.e]+),\s*([-\d.e]+)/);if(m)a+=Math.atan2(+m[2],+m[1])*180/Math.PI;n=n.parentElement}return a};
 const rgb=s=>{const m=(s||'').match(/[\d.]+/g);return m?m.slice(0,4).map(Number):null};
 const direct=el=>[...el.childNodes].some(n=>n.nodeType===3&&n.textContent.trim());
 const skip=el=>el.closest('#b-scheme')||['SCRIPT','STYLE','IMG','SVG','CANVAS'].includes(el.tagName);
 for(const el of document.querySelectorAll('#frame *,#txt *')){
  if(skip(el)||!direct(el))continue; const r=el.getBoundingClientRect(),cs=getComputedStyle(el);
  if(r.width<2||r.height<2||cs.visibility==='hidden'||+cs.opacity===0)continue;
  const own=[...el.childNodes].filter(n=>n.nodeType===3).map(n=>n.textContent).join('').trim();
  if(!own)continue;
  texts.push({text:own,x:r.left,y:r.top,w:r.width,h:r.height,rot:angle(el),size:parseFloat(cs.fontSize),
   line:parseFloat(cs.lineHeight)||parseFloat(cs.fontSize)*1.25,weight:parseInt(cs.fontWeight)||400,
   color:rgb(cs.color),align:cs.textAlign,family:cs.fontFamily.split(',')[0].replace(/["']/g,'').trim(),
   spacing:parseFloat(cs.letterSpacing)||0,italic:cs.fontStyle==='italic'});
 }
 for(const el of document.querySelectorAll('#stk img,#txt img')){
  const r=el.getBoundingClientRect(),cs=getComputedStyle(el); if(r.width<2||r.height<2||+cs.opacity===0)continue;
  images.push({src:el.currentSrc,x:r.left,y:r.top,w:r.width,h:r.height,rot:angle(el)});
 }
 const candidates=new Set([...document.querySelectorAll('#frame .gl,#frame .xm,#prog,#prog i,#txt .chip,#txt .tc,#txt .row,#txt .it,#txt .st-i,#txt .cm,#txt .sp,#txt .pc,#txt .grow,#txt .rule')]);
 for(const el of candidates){const r=el.getBoundingClientRect(),cs=getComputedStyle(el);if(r.width<1||r.height<1)continue;
  if(el.classList.contains('xm')){shapes.push({kind:'cross',x:r.left,y:r.top,w:r.width,h:r.height,color:[0,0,0],alpha:.42});continue}
  const bg=rgb(cs.backgroundColor),bc=rgb(cs.borderTopColor); const bw=Math.max(parseFloat(cs.borderTopWidth)||0,parseFloat(cs.borderRightWidth)||0,parseFloat(cs.borderBottomWidth)||0,parseFloat(cs.borderLeftWidth)||0);
  shapes.push({kind:'rect',x:r.left,y:r.top,w:r.width,h:r.height,rot:angle(el),bg:bg,border:bc,bw:bw,alpha:parseFloat(cs.opacity)||1,radius:parseFloat(cs.borderRadius)||0});
 }
 const rootBg=rgb(getComputedStyle(document.body).backgroundColor);
 const api=window.__PPTX_EXPORT__;
 const current=+(new URLSearchParams(location.search).get('p')||0);
 const exportMeta={pageCount:api?.pageCount?.()||document.querySelectorAll('[data-slide],.slide').length||1,
   objects:api?.objectsForPage?.(current)||[]};
 const data={texts,images,shapes,rootBg,exportMeta}; const pre=document.createElement('pre');pre.id='__hxout';
 pre.textContent=btoa(unescape(encodeURIComponent(JSON.stringify(data))));document.body.appendChild(pre);document.title='HXREADY';
})();
</script>
"""

class Quiet(http.server.SimpleHTTPRequestHandler):
    def log_message(self, *_): pass

@contextlib.contextmanager
def server(root):
    handler=lambda *a,**k: Quiet(*a,directory=str(root),**k)
    with socketserver.TCPServer(("127.0.0.1",0),handler) as httpd:
        t=threading.Thread(target=httpd.serve_forever,daemon=True);t.start()
        try: yield httpd.server_address[1]
        finally: httpd.shutdown();t.join()

def chrome(args, timeout=35):
    cmd=[CHROME,"--headless=new","--hide-scrollbars","--disable-background-networking","--disable-default-apps","--disable-extensions","--no-first-run","--no-default-browser-check",*args]
    try: return subprocess.run(cmd,capture_output=True,timeout=timeout)
    except subprocess.TimeoutExpired as e: return subprocess.CompletedProcess(cmd,124,e.stdout or b'',e.stderr or b'')

def rgba(values, default=(0,0,0,0)):
    if not values:return default
    vals=list(values)+[1];return tuple(vals[:4])

def main():
    ap=argparse.ArgumentParser();ap.add_argument('deck');ap.add_argument('-o','--out');ap.add_argument('--scheme',choices=['base','plus','bold'],default='plus');args=ap.parse_args()
    from pptx import Presentation
    from pptx.util import Emu,Pt
    from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    deck=Path(args.deck).resolve()
    # Serve from a stable common project root while still allowing nested/custom decks.
    cwd=Path.cwd().resolve()
    root=cwd if deck.is_relative_to(cwd) else next((p for p in deck.parents if (p/'scripts'/'export-haoqi-pptx-editable.py').exists()),deck.parent)
    rel=deck.relative_to(root)
    html=deck.read_text(); injected=html.replace('</body>',COLLECT+'</body>'); probe=deck.with_name('.export-probe.html');probe.write_text(injected)
    td=Path(tempfile.mkdtemp(prefix='haoqi-pptx-')); pages=[]; backgrounds=[]; objects=[]
    try:
      with server(root) as port:
       probe_rel=probe.relative_to(root)
       meta_url=f'http://127.0.0.1:{port}/{probe_rel}?p=0&still=1&s={args.scheme}'
       first=chrome([f'--window-size={VW},{VH}','--virtual-time-budget=7000','--dump-dom',meta_url])
       mm=re.search(rb'id="__hxout">([A-Za-z0-9+/=]+)<',first.stdout or b'')
       if not mm: raise SystemExit('ERROR: deck metadata extraction failed; ensure the HTML closes with </body>')
       first_data=json.loads(base64.b64decode(mm.group(1)));n=int(first_data.get('exportMeta',{}).get('pageCount',0))
       if n < 1: raise SystemExit('ERROR: deck reports zero pages')
       print(f'  ✓ discovered {n} slide(s) from runtime metadata')
       for i in range(n):
        base=f'http://127.0.0.1:{port}/{probe_rel}?p={i}&still=1&s={args.scheme}'
        r=chrome([f'--window-size={VW},{VH}','--virtual-time-budget=7000','--dump-dom',base])
        m=re.search(rb'id="__hxout">([A-Za-z0-9+/=]+)<',r.stdout or b'')
        if not m: raise SystemExit(f'ERROR: page {i+1} DOM extraction failed')
        pages.append(first_data if i==0 else json.loads(base64.b64decode(m.group(1))))
        reported=int(pages[-1].get('exportMeta',{}).get('pageCount',n))
        if reported != n: raise SystemExit(f'ERROR: page count changed during export ({n} → {reported})')
        bg=td/f'bg-{i+1:02d}.png'; url=f'http://127.0.0.1:{port}/{rel}?p={i}&still=1&s={args.scheme}&export=background'
        chrome([f'--window-size={VW},{VH}','--virtual-time-budget=6500',f'--screenshot={bg}',url])
        if not bg.exists(): raise SystemExit(f'ERROR: page {i+1} WebGL background failed')
        try:
          from PIL import Image
          jpg=bg.with_suffix('.jpg');Image.open(bg).convert('RGB').save(jpg,'JPEG',quality=88);bg.unlink();bg=jpg
        except ImportError: pass
        backgrounds.append(bg)
        page_objects=[]
        for oi,obj_id in enumerate(pages[-1].get('exportMeta',{}).get('objects',[])):
          raw=td/f'obj-{i+1:03d}-{oi:02d}-raw.png'; qobj=urllib.parse.quote(str(obj_id),safe='')
          objurl=f'http://127.0.0.1:{port}/{rel}?p={i}&still=1&s={args.scheme}&export=object&object={qobj}'
          chrome([f'--window-size={VW},{VH}','--virtual-time-budget=6500',f'--screenshot={raw}',objurl])
          if raw.exists():
           try:
            from PIL import Image
           except ImportError: raise SystemExit('ERROR: Pillow is required when exporting independent 3D objects')
           im=Image.open(raw).convert('RGBA'); box=im.getchannel('A').getbbox()
           if box:
            obj=td/f'obj-{i+1:03d}-{oi:02d}.png';im.crop(box).save(obj);page_objects.append((obj,box,str(obj_id)))
        objects.append(page_objects);print(f'  ✓ capture {i+1:03d}/{n}: {len(pages[-1]["texts"])} text · {len(pages[-1]["images"])} DOM image · {len(page_objects)} 3D object · {len(pages[-1]["shapes"])} shape')
    finally: probe.unlink(missing_ok=True)

    prs=Presentation();prs.slide_width=Emu(12192000);prs.slide_height=Emu(6858000);blank=prs.slide_layouts[6]
    aligns={'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT,'justify':PP_ALIGN.JUSTIFY}
    fontmap={'-apple-system':'PingFang SC','Helvetica Neue':'Arial','PingFang SC':'PingFang SC','ui-monospace':'Menlo'}
    counts={'background':0,'image':0,'text':0,'shape':0}
    for si,data in enumerate(pages):
      slide=prs.slides.add_slide(blank)
      bgc=rgba(data.get('rootBg'),(191,221,240,1));fill=slide.background.fill;fill.solid();fill.fore_color.rgb=RGBColor(*map(int,bgc[:3]))
      slide.shapes.add_picture(str(backgrounds[si]),0,0,width=prs.slide_width,height=prs.slide_height);counts['background']+=1
      for obj,box,obj_id in objects[si]:
       x0,y0,x1,y1=box
       pic=slide.shapes.add_picture(str(obj),Emu(int(x0*EMU_PER_PX)),Emu(int(y0*EMU_PER_PX)),width=Emu(int((x1-x0)*EMU_PER_PX)),height=Emu(int((y1-y0)*EMU_PER_PX)))
       pic.name=f'3D: {obj_id}';counts['image']+=1
      for sh in data['shapes']:
       if sh['kind']=='cross':
        for x,y,w,h in [(sh['x']+sh['w']/2,sh['y'],1,sh['h']),(sh['x'],sh['y']+sh['h']/2,sh['w'],1)]:
         z=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Emu(int(x*EMU_PER_PX)),Emu(int(y*EMU_PER_PX)),Emu(max(1,int(w*EMU_PER_PX))),Emu(max(1,int(h*EMU_PER_PX))));z.line.fill.background();z.fill.solid();z.fill.fore_color.rgb=RGBColor(0,0,0);z.fill.transparency=58;counts['shape']+=1
        continue
       bg=rgba(sh.get('bg')); border=rgba(sh.get('border')); has_bg=bg[3]>.01; has_border=sh.get('bw',0)>.1 and border[3]>.01
       if not has_bg and not has_border:continue
       z=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if sh.get('radius',0)>3 else MSO_SHAPE.RECTANGLE,Emu(int(sh['x']*EMU_PER_PX)),Emu(int(sh['y']*EMU_PER_PX)),Emu(int(sh['w']*EMU_PER_PX)),Emu(int(sh['h']*EMU_PER_PX)));z.rotation=sh.get('rot',0)
       if has_bg:z.fill.solid();z.fill.fore_color.rgb=RGBColor(*map(int,bg[:3]));z.fill.transparency=max(0,min(100,int((1-bg[3]*sh.get('alpha',1))*100)))
       else:z.fill.background()
       if has_border:z.line.color.rgb=RGBColor(*map(int,border[:3]));z.line.width=Pt(max(.4,sh['bw']*PT_PER_PX));z.line.transparency=max(0,min(100,int((1-border[3])*100)))
       else:z.line.fill.background()
       counts['shape']+=1
      for im in data['images']:
       u=urllib.parse.urlparse(im['src']);src=Path(urllib.parse.unquote(u.path));
       # URL path is relative to HTTP root
       src=root/str(src).lstrip('/')
       if not src.exists():raise SystemExit(f'ERROR missing image {src}')
       pic=slide.shapes.add_picture(str(src),Emu(int(im['x']*EMU_PER_PX)),Emu(int(im['y']*EMU_PER_PX)),width=Emu(int(im['w']*EMU_PER_PX)),height=Emu(int(im['h']*EMU_PER_PX)));pic.rotation=im.get('rot',0);counts['image']+=1
      for it in data['texts']:
       w=max(it['w']*1.10,12);h=max(it['h']*1.18,8);box=slide.shapes.add_textbox(Emu(int((it['x']-w*.025)*EMU_PER_PX)),Emu(int((it['y']-h*.07)*EMU_PER_PX)),Emu(int(w*EMU_PER_PX)),Emu(int(h*EMU_PER_PX)));box.rotation=it.get('rot',0)
       tf=box.text_frame;tf.clear();tf.word_wrap=True;tf.vertical_anchor=MSO_ANCHOR.MIDDLE;tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
       for li,line in enumerate(it['text'].split('\n')):
        p=tf.paragraphs[0] if li==0 else tf.add_paragraph();p.alignment=aligns.get(it.get('align'),PP_ALIGN.LEFT);p.line_spacing=max(.9,it['line']/it['size']);run=p.add_run();run.text=line;f=run.font;f.size=Pt(round(it['size']*PT_PER_PX,1));f.bold=it['weight']>=600;f.italic=it.get('italic',False);c=rgba(it['color']);f.color.rgb=RGBColor(*map(int,c[:3]));f.name=fontmap.get(it['family'],it['family'])
       counts['text']+=1
    out=Path(args.out).expanduser().resolve() if args.out else deck.with_name(deck.stem+'-editable.pptx');out.parent.mkdir(parents=True,exist_ok=True);prs.save(out)
    print(f'  ✓ {out} ({out.stat().st_size/1024/1024:.1f}MB)')
    print('  ✓ objects:',json.dumps(counts,ensure_ascii=False))

if __name__=='__main__':main()
